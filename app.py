import streamlit as st
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer
from google import genai
from dotenv import load_dotenv
import os
import fitz  # PyMuPDF
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from docx import Document
from io import BytesIO
import logging

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('rag_app.log')
    ]
)
logger = logging.getLogger(__name__)

load_dotenv()

if 'documents_loaded' not in st.session_state:
    st.session_state.documents_loaded = False
if 'collection' not in st.session_state:
    st.session_state.collection = None
if 'embedding_model' not in st.session_state:
    st.session_state.embedding_model = None

@st.cache_resource
def initialize_embeddings():
    """Load Hugging Face embedding model"""
    model = SentenceTransformer('BAAI/bge-small-en-v1.5')
    return model

def initialize_chroma():
    client = chromadb.PersistentClient(path="./chroma_db")
    collection = client.get_or_create_collection(
        name="documents",
        metadata={"hnsw:space": "cosine"}
    )
    return collection

def load_documents(uploaded_files):
    documents = []
    for uploaded_file in uploaded_files:
        file_extension = uploaded_file.name.split('.')[-1].lower() if '.' in uploaded_file.name else ''
        
        if uploaded_file.type == "application/pdf" or file_extension == "pdf":
            pdf_doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            text = ""
            for page in pdf_doc:
                text += page.get_text()
            pdf_doc.close()
            documents.append(text)
        elif uploaded_file.type == "text/plain" or file_extension == "txt":
            text = str(uploaded_file.read(), "utf-8")
            documents.append(text)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or file_extension == "docx":
            # DOCX file
            doc = Document(BytesIO(uploaded_file.read()))
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            documents.append(text)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or file_extension == "pptx":
            # PPTX file
            prs = Presentation(BytesIO(uploaded_file.read()))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            documents.append(text)
    return documents

def chunk_text(texts):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        length_function=len,
    )
    chunks = []
    for text in texts:
        text_chunks = text_splitter.split_text(text)
        chunks.extend(text_chunks)
    return chunks

def add_documents_to_db(chunks, embedding_model, collection):
    """Store document chunks in ChromaDB with embeddings"""
    if not chunks:
        return
    
    embeddings = embedding_model.encode(chunks).tolist()
    
    ids = [f"chunk_{i}" for i in range(len(chunks))]
    
    collection.add(
        embeddings=embeddings,
        documents=chunks,
        ids=ids
    )

def query_rag(question, embedding_model, collection, api_key):
    logger.info(f"QUERY RECEIVED: {question}")
    
    # Embed the question
    question_embedding = embedding_model.encode([question]).tolist()[0]
    logger.info("Question embedded successfully")
    
    # Retrieve more chunks for better context
    results = collection.query(
        query_embeddings=[question_embedding],
        n_results=5
    )
    
    retrieved_docs = results['documents'][0] if results['documents'] else []
    logger.info(f"Retrieved {len(retrieved_docs)} relevant chunks from ChromaDB")
    
    # Log retrieved chunks
    for i, doc in enumerate(retrieved_docs, 1):
        logger.info(f"Chunk {i} (length: {len(doc)}")
    
    context = "\n\n".join(retrieved_docs)
    logger.info(f"Total context length: {len(context)} characters")
    
    # Create enhanced prompt for more detailed answers
    prompt = f"""You are a helpful assistant that provides detailed and comprehensive answers based on the given context.

        Context from documents:
        {context}

        Question: {question}

        Instructions:
        - Provide a detailed, comprehensive answer based on the context above
        - Include relevant details, examples, and explanations from the context
        - If the answer is not fully available in the context, clearly state what information is available and what is missing
        - Structure your answer clearly with proper formatting
        - Be thorough and informative

        Answer:"""
    
    logger.info("CALLING GEMINI API")
    logger.info(f"Prompt length: {len(prompt)} characters")
    logger.info(f"PROMPT SENT TO LLM:\n{prompt}")
    logger.info("-" * 80)
    
    try:
        client = genai.Client(api_key=api_key)
        response = client.models.generate_content(
            model='gemini-2.5-flash-lite',
            contents=prompt
        )
        
        answer = response.text
        logger.info("RESPONSE RECEIVED FROM GEMINI API")
        logger.info(f"Answer length: {len(answer)} characters")
        logger.info(f"ANSWER:\n{answer}")
        
        return answer
    except Exception as e:
        logger.error(f"ERROR in Gemini API call: {str(e)}")
        logger.error(f"Error type: {type(e).__name__}")
        raise

st.title("Basic RAG Application")
st.write("Upload documents and ask questions!")

with st.sidebar:
    st.header("Upload Documents")
    uploaded_files = st.file_uploader(
        "Choose files",
        type=['pdf', 'txt', 'docx', 'pptx'],
        accept_multiple_files=True
    )
    
    if st.button("Process Documents"):
        if uploaded_files:
            with st.spinner("Processing documents..."):
                if st.session_state.embedding_model is None:
                    st.session_state.embedding_model = initialize_embeddings()
                
                if st.session_state.collection is None:
                    st.session_state.collection = initialize_chroma()
                
                documents = load_documents(uploaded_files)
                logger.info(f"Loaded {len(documents)} document(s) from {len(uploaded_files)} file(s)")
                
                chunks = chunk_text(documents)
                logger.info(f"Created {len(chunks)} chunks from documents")
                
                add_documents_to_db(
                    chunks,
                    st.session_state.embedding_model,
                    st.session_state.collection
                )
                logger.info(f"Added {len(chunks)} chunks to ChromaDB")
                
                st.session_state.documents_loaded = True
                st.success(f"Processed {len(chunks)} chunks from {len(uploaded_files)} file(s)!")
        else:
            st.warning("Please upload at least one file.")

# Main area
if st.session_state.documents_loaded:
    st.success("Documents loaded! You can now ask questions.")
    
    question = st.text_input("Ask a question:")
    
    if question:
        api_key = os.getenv("GEMINI_API_KEY")
        if api_key:
            with st.spinner("Generating answer..."):
                try:
                    answer = query_rag(
                        question,
                        st.session_state.embedding_model,
                        st.session_state.collection,
                        api_key
                    )
                    st.write("**Answer:**")
                    st.write(answer)
                except Exception as e:
                    logger.error(f"Error in query_rag: {str(e)}", exc_info=True)
                    st.error(f"Error: {str(e)}")
        else:
            st.error("GEMINI_API_KEY not found in .env file")
else:
    st.info("Please upload and process documents first.")
